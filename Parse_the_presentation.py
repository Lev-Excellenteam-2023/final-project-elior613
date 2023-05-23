from pptx import Presentation
from io import BytesIO
from PIL import Image

def extract_content_from_shape(shape):
    """
    Extracts content from a shape based on its type (image or table).
    """
    content = ""

    if shape.shape_type == 13:  # Check if the shape is an image (shape_type=13)
        content += extract_image_content(shape)+" "
    elif shape.shape_type == 19:  # Check if the shape is a table (shape_type=19)
        content += extract_table_content(shape)+" "
    else:
        content += extract_text_from_shape(shape)+" "

    return content.strip()

def extract_image_content(shape):
    """
    Extracts and returns the content of an embedded image shape.
    """
    image = shape.image
    image_data = image.blob
    image_file = BytesIO(image_data)
    image_pil = Image.open(image_file)
    # Process the image and extract relevant content (e.g., using computer vision techniques)
    # Replace the following line with your custom image processing code
    content = f"[IMAGE CONTENT: {image.filename}]"
    return content

def extract_table_content(shape):
    """
    Extracts the contents of a table shape.
    """
    content = ""
    table = shape.table

    for row in table.rows:
        for cell in row.cells:
            cell_text = extract_text_from_shape(cell)
            content += cell_text + " "

    return content.strip()

def extract_text_from_shape(shape, title_text=None):
    """
    Extracts text from a shape's text frame paragraphs and runs, excluding the title text.
    """
    text = ""
    if shape.has_text_frame and not shape.has_table:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text not in title_text:
                    text += run.text + " "
    return text.strip()


def extract_content_from_slide(slide):
    """
    Extracts the title and content from a slide, including images and tables.
    """
    title = slide.shapes.title.text if slide.shapes.title else "NONE"
    content = ""

    for shape in slide.shapes:
        content += extract_text_from_shape(shape, title)+" "

    return title, content


def process_presentation(presentation_path):
    """
    Processes a PowerPoint presentation and returns a list of dictionaries containing slide information.
    """
    presentation = Presentation(presentation_path)
    slides_data = []

    for slide_num, slide in enumerate(presentation.slides, start=1):
        title, content = extract_content_from_slide(slide)
        slide_data = {
            "slide_number": slide_num,
            "title": title,
            "content": content,
        }
        slides_data.append(slide_data)

    return slides_data

def print_slide_content(slides_data):
    """
    Prints the slide number, title, content, and image subjects of slides.
    """
    for slide_data in slides_data:
        print(f"Slide {slide_data['slide_number']}")
        print(f"Title: {slide_data['title']}")
        print(f"Content: {slide_data['content']}")
        print()

# Usage example
presentation_path = "C:\machine learning.pptx"
slides_data = process_presentation(presentation_path)
print_slide_content(slides_data)
